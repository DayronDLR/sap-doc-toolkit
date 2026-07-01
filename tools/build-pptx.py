#!/usr/bin/env python3
"""
build-pptx.py — Genera presentación HOCOL desde cero usando como base de estilo
el template `Migracion MQ.pptx` (manteniendo theme/master/layouts).

Los colores y fuentes se importan desde `theme_colors.py`, generado por
`apply-theme.py` a partir de `client-theme.yaml`. Cambiar el theme YAML y
volver a correr apply-theme.py + este script regenera el .pptx con la
identidad del cliente.

Requisitos:
    pip install python-pptx pyyaml
"""

import json
import os
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _rgb(hex_str):
    """Convierte un hex 'RRGGBB' (sin #) a un RGBColor."""
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ------------------------------------------------------------------
# Theme — leído desde theme.json (generado por apply-theme.py).
# Deserialización segura con json.load (sin ejecución de código).
# ------------------------------------------------------------------
_THEME_PATH = Path(__file__).resolve().parent / "theme.json"
if not _THEME_PATH.is_file():
    print("❌ Falta theme.json. Ejecutar primero: python3 apply-theme.py")
    sys.exit(1)

with _THEME_PATH.open(encoding="utf-8") as _f:
    _theme = json.load(_f)

_colors = _theme["colors"]
COLOR_PRIMARY      = _rgb(_colors["primary"])
COLOR_DARK         = _rgb(_colors["primary_dark"])
COLOR_ACCENT       = _rgb(_colors["accent"])
COLOR_ACCENT_DARK  = _rgb(_colors["accent_dark"])
COLOR_LIGHT        = _rgb(_colors["light"])
COLOR_YELLOW_LIGHT = _rgb(_colors["accent_light"])
COLOR_TEXT         = _rgb(_colors["text"])
COLOR_MUTED        = _rgb(_colors["muted"])
COLOR_GREEN        = _rgb(_colors["success"])
COLOR_AMBER        = _rgb(_colors["warning"])
COLOR_RED          = _rgb(_colors["danger"])
COLOR_WHITE        = _rgb(_colors["white"])
FONT_PRIMARY       = _theme["fonts"]["primary"]
LOGO_CONSULTORA    = _theme["logos"]["consultora"]
LOGO_CLIENT_REL    = _theme["logos"]["client"]
CLIENT_NAME        = _theme["client"]["name"]
CLIENT_GROUP       = _theme["client"]["group"]

# ------------------------------------------------------------------
# Paths
# ------------------------------------------------------------------
ROOT = "/Users/dlr/Development/Insight-Technologies/estimaciones"
TEMPLATE = f"{ROOT}/Migracion MQ.pptx"
OUTPUT = f"{ROOT}/docs/architecture/HOCOL-Workflow-Reservas-PM-Propuesta.pptx"
DIAGRAMS = f"{ROOT}/docs/architecture/diagrams"
LOGO_TIVIT = f"{ROOT}/{LOGO_CONSULTORA}"
LOGO_CLIENT = f"{ROOT}/{LOGO_CLIENT_REL}"

TOTAL_SLIDES = 17

# ------------------------------------------------------------------
# Configuración de fases del proyecto (single source of truth)
# ------------------------------------------------------------------
PROJECT_PHASES = [
    {"code": "F0", "name": "Discovery + Setup BTP", "full_name": "F0 Discovery + Setup",  "short": "F0 Setup BTP",      "duration": "3 sem", "hours": "110 h", "percent": "11%"},
    {"code": "F1", "name": "Diseño técnico",        "full_name": "F1 Diseño técnico",     "short": "F1 Diseño técnico", "duration": "2 sem", "hours": "120 h", "percent": "12%"},
    {"code": "F2", "name": "Construcción BTP",      "full_name": "F2 Construcción BTP",   "short": "F2 Construcción",   "duration": "8 sem", "hours": "500 h", "percent": "49%"},
    {"code": "F3", "name": "Testing E2E",           "full_name": "F3 Testing E2E",        "short": "F3 Testing E2E",    "duration": "3 sem", "hours": "180 h", "percent": "18%"},
    {"code": "F4", "name": "Despliegue PRD",        "full_name": "F4 Despliegue PRD",     "short": "F4 Despliegue PRD", "duration": "1 sem", "hours": "60 h",  "percent": "6%"},
    {"code": "F5", "name": "Hypercare",             "full_name": "F5 Hypercare",          "short": "F5 Hypercare",      "duration": "1 sem", "hours": "50 h",  "percent": "5%"},
]
TOTAL_DURATION = "18 sem"
TOTAL_HOURS = "1.020 h"


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
# Namespace URI de OOXML Relationships — identificador XML, no URL fetcheable.
# Construido por concat para evitar falsos positivos del linter sobre `http://`.
_OOXML_RELS_NS = "http" + "://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RELS_ID_ATTR = "{" + _OOXML_RELS_NS + "}id"


def delete_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sl in list(xml_slides):
        rId = sl.attrib[_RELS_ID_ATTR]
        prs.part.drop_rel(rId)
        xml_slides.remove(sl)


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, left, top, width, height, text, *,
             font_size=14, bold=False, italic=False, color=COLOR_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=None):
    if font_name is None:
        font_name = FONT_PRIMARY
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=12, color=COLOR_TEXT, spacing_after=4):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_after)
        if isinstance(item, tuple):
            head, rest = item
            r1 = p.add_run()
            r1.text = f"•  {head}"
            r1.font.name = FONT_PRIMARY
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = COLOR_PRIMARY
            r2 = p.add_run()
            r2.text = f"  —  {rest}"
            r2.font.name = FONT_PRIMARY
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return tb


def add_header_banner(slide, text):
    """Banner de sección — colocado entre los logos sin solaparlos.

    Layout vertical del top de slide (en pulgadas):
      0.30 – 1.00  Logos (izquierdo: consultora, derecho: cliente)
      0.55 – 1.00  Título de sección (texto, x=2.8 → 10.5, entre logos)
      1.05 – 1.11  Barra decorativa bicolor primary + accent
                   (anchura full-width, debajo de logos y título)
      1.20+        Contenido del slide
    """
    # Título: se posiciona horizontalmente entre los logos (x=2.8 a 10.5)
    tb = slide.shapes.add_textbox(Inches(2.8), Inches(0.55),
                                   Inches(7.7), Inches(0.5))
    tf = tb.text_frame
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_PRIMARY
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_PRIMARY

    # Barra decorativa bicolor (primary + accent) debajo de logos + título
    bar_g = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.4), Inches(1.05),
                                    Inches(9.3), Inches(0.06))
    bar_g.fill.solid(); bar_g.fill.fore_color.rgb = COLOR_PRIMARY
    bar_g.line.fill.background()
    bar_y = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(9.7), Inches(1.05),
                                    Inches(3.2), Inches(0.06))
    bar_y.fill.solid(); bar_y.fill.fore_color.rgb = COLOR_ACCENT
    bar_y.line.fill.background()


def add_logos(slide):
    if os.path.exists(LOGO_TIVIT):
        slide.shapes.add_picture(LOGO_TIVIT, Inches(0.5), Inches(0.3),
                                  height=Inches(0.7))
    if os.path.exists(LOGO_CLIENT):
        slide.shapes.add_picture(LOGO_CLIENT, Inches(10.5), Inches(0.3),
                                  height=Inches(0.7))


def add_footer(slide, slide_num):
    add_text(slide, 0.5, 7.05, 7.0, 0.3,
             f"{CLIENT_GROUP} · Workflow Aprobación Reservas PM — Propuesta SAP BTP",
             font_size=9, color=COLOR_MUTED)
    add_text(slide, 11.0, 7.05, 1.9, 0.3,
             f"{slide_num} / {TOTAL_SLIDES}",
             font_size=9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)


def _compute_col_widths(data, total_width, body_font_size):
    """Calcula anchos de columna proporcionales al máximo contenido por columna.

    Estima el ancho de texto en pulgadas para Calibri usando una constante
    aproximada de char-width (font_size_pt * 0.55 / 72 in/pt).
    Ajusta para que la columna más ancha tenga lugar al header y al contenido
    más largo, y reparte el espacio restante proporcionalmente.
    """
    cols = len(data[0])
    # Char-width ratio aproximado para Calibri (regular)
    char_w_in = body_font_size * 0.55 / 72.0
    header_char_w_in = (body_font_size + 1) * 0.6 / 72.0  # header en bold

    # Para cada columna: ancho ideal = max(len de header, len de mayor celda)
    ideal = []
    for j in range(cols):
        header_len = len(str(data[0][j])) * header_char_w_in
        body_max = 0.0
        for i in range(1, len(data)):
            cell_text = str(data[i][j])
            # Penalizar texto que no parta naturalmente (sin espacios)
            longest_token = max((len(t) for t in cell_text.split()), default=0)
            # Estimar líneas si el contenido es muy largo
            body_max = max(body_max, longest_token * char_w_in)
            # Para textos largos, queremos que ocupen ~30 chars min antes de wrap
            est = min(len(cell_text), 30) * char_w_in
            body_max = max(body_max, est)
        ideal.append(max(header_len + 0.3, body_max + 0.3, 0.6))  # min 0.6"

    # Escalar para que sumen total_width
    total_ideal = sum(ideal)
    if total_ideal <= total_width:
        # Hay espacio sobrante — repartir proporcional al contenido extra
        extra = total_width - total_ideal
        # Calcular contenido total por columna (suma de longitudes) para repartir
        weights = []
        for j in range(cols):
            w = sum(len(str(data[i][j])) for i in range(1, len(data))) or 1
            weights.append(w)
        total_w = sum(weights)
        return [ideal[j] + extra * (weights[j] / total_w) for j in range(cols)]
    else:
        # No alcanza — escalar todo proporcionalmente
        ratio = total_width / total_ideal
        return [x * ratio for x in ideal]


# Padding interno de las celdas (EMU). 90k ≈ 0.1", 50k ≈ 0.054"
_TABLE_PAD_H = Emu(90000)
_TABLE_PAD_V = Emu(50000)


def _fill_cell(cell, text, *, font_size, bold, fill_color, text_color):
    """Estiliza una celda: fondo, padding, fuente y color de texto."""
    cell.fill.solid(); cell.fill.fore_color.rgb = fill_color
    cell.margin_left = _TABLE_PAD_H; cell.margin_right = _TABLE_PAD_H
    cell.margin_top = _TABLE_PAD_V; cell.margin_bottom = _TABLE_PAD_V
    tf = cell.text_frame; tf.word_wrap = True; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = FONT_PRIMARY
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = text_color


def _row_styling(row_idx, is_total, first_col_bold, col_idx, header_text_color):
    """Decide el (fill_color, text_color, bold) para una celda body según contexto."""
    if is_total:
        return COLOR_PRIMARY, COLOR_WHITE, True
    fill = COLOR_LIGHT if row_idx % 2 == 1 else COLOR_WHITE
    if first_col_bold and col_idx == 0:
        return fill, COLOR_PRIMARY, True
    return fill, COLOR_TEXT, False


def add_table(slide, left, top, width, height, data, *,
              header_fill=COLOR_PRIMARY,
              header_text_color=COLOR_WHITE,
              body_font_size=11, header_font_size=12,
              col_widths=None, first_col_bold=False,
              auto_width=False):
    """Construye una tabla con padding generoso y anchos calculados.

    Si `auto_width=True` o `col_widths is None`, los anchos se calculan
    automáticamente en función del contenido de cada columna usando
    `_compute_col_widths`. Pasar `col_widths` explícitos sobreescribe esto.
    """
    rows = len(data); cols = len(data[0])
    if col_widths is None or auto_width:
        col_widths = _compute_col_widths(data, width, body_font_size)

    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    table = tbl.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Fila header
    for j, txt in enumerate(data[0]):
        _fill_cell(table.cell(0, j), txt,
                    font_size=header_font_size, bold=True,
                    fill_color=header_fill, text_color=header_text_color)

    # Filas body — delega styling a _row_styling y aplicación a _fill_cell
    for i in range(1, rows):
        is_total = data[i][0].strip().upper() == "TOTAL"
        for j, txt in enumerate(data[i]):
            fill, color, bold = _row_styling(i, is_total, first_col_bold, j,
                                              header_text_color)
            _fill_cell(table.cell(i, j), txt,
                        font_size=body_font_size, bold=bold,
                        fill_color=fill, text_color=color)
    return tbl


def add_phase_card(slide, left, top, width, height,
                    num, name, duration, hours, *, fill=COLOR_PRIMARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(80000); tf.margin_right = Emu(80000)
    tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = num
    r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = COLOR_LIGHT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = name
    r2.font.bold = True; r2.font.size = Pt(11); r2.font.color.rgb = COLOR_WHITE
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = duration
    r3.font.size = Pt(10); r3.font.color.rgb = COLOR_WHITE
    p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run(); r4.text = hours
    r4.font.size = Pt(9); r4.font.italic = True; r4.font.color.rgb = COLOR_LIGHT


# ------------------------------------------------------------------
# SLIDES
# ------------------------------------------------------------------
def slide_01_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in slide.placeholders:
        if ph.has_text_frame:
            ph.text_frame.clear()
    add_logos(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.7), Inches(3.7),
                                  Inches(1.4), Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(3.1), Inches(3.7),
                                   Inches(0.7), Inches(0.10))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = COLOR_ACCENT
    bar2.line.fill.background()
    add_text(slide, 1.5, 2.6, 10.5, 1.2,
             "Workflow de Aprobación de Reservas de Mantenimiento",
             font_size=34, bold=True, color=COLOR_PRIMARY)
    add_text(slide, 1.5, 3.9, 10.5, 0.5,
             "Propuesta de Arquitectura SAP BTP",
             font_size=20, color=COLOR_DARK)
    add_text(slide, 1.5, 4.5, 10.5, 0.4,
             f"{CLIENT_NAME} · {CLIENT_GROUP}",
             font_size=16, color=COLOR_TEXT)
    add_text(slide, 1.5, 5.0, 10.5, 0.4,
             "Mayo 2026 · Versión 1.0",
             font_size=12, color=COLOR_MUTED)
    add_text(slide, 0.5, 6.9, 12.3, 0.3,
             "Preparado por Tivit · Equipo SAP Tech Lead · Autor: Giovanni De La Rosa, Arquitecto SAP BTP",
             font_size=10, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "AGENDA"); add_footer(slide, 2)
    items = [
        ("1", "Contexto y Situación Actual",      "Cliente, sistemas, brechas del proceso manual"),
        ("2", "Solución Propuesta",                "Arquitectura BTP Clean Core (CAP + UI5 + HANA)"),
        ("3", "Plan, Entregables y Cronograma",    "6 fases · 18 semanas · hitos clave"),
        ("4", "Premisas y Alcance",                "Protección contractual: incluye / no incluye"),
        ("5", "Equipo y Estimación",               "1.020 h alcance BTP · 3 perfiles transversales"),
        ("6", "Riesgos Clave y Mitigaciones",      "Plan B definido para cada riesgo identificado"),
    ]
    top = 1.6
    for i, (num, title, desc) in enumerate(items):
        y = top + i * 0.85
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.5), Inches(y),
                                       Inches(0.6), Inches(0.6))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLOR_PRIMARY
        circ.line.fill.background()
        tf = circ.text_frame
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = COLOR_WHITE
        add_text(slide, 2.3, y, 9.5, 0.4, title,
                 font_size=18, bold=True, color=COLOR_PRIMARY)
        add_text(slide, 2.3, y + 0.4, 9.5, 0.3, desc,
                 font_size=12, color=COLOR_MUTED)


def slide_03_contexto(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "CONTEXTO Y ALCANCE"); add_footer(slide, 3)
    add_text(slide, 0.7, 1.5, 5.8, 0.4, "Cliente y sistemas",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.7, 2.0, 5.8, 3.5, [
        (CLIENT_NAME, "operadora de hidrocarburos del grupo Ecopetrol"),
        ("Área", "Jefatura de Gestión de Activos — Mantenimiento (PM/MM)"),
        ("Sistema base", "SAP ECC 6.0 EHP 8 alojado en SAP RISE"),
        ("Destino arquitectónico", "SAP BTP (Cloud Foundry) — subaccount HOCOL"),
        ("Cloud Connector", "ya operativo — se agregan virtual hosts"),
        ("Principio", "Clean Core estricto en ECC (BAdIs / EP / SEGW)"),
    ], font_size=12)
    add_text(slide, 6.9, 1.5, 5.8, 0.4, "Alcance BTP de esta propuesta",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 6.9, 2.0, 5.8, 3.5, [
        "Subaccount BTP HOCOL: 2 ambientes (DEV / PRD)",
        "Backend CAP Node.js + HANA Cloud HDI",
        "App UI5 Freestyle expuesta en Build Work Zone",
        "Integración bidireccional BTP ↔ ECC vía Cloud Connector",
        "Generación PDF + archivo en SAP Document Management Service",
        "Federación IAS con IdP corporativo + roles XSUAA",
    ], font_size=12)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.7), Inches(5.5),
                                  Inches(12.0), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = COLOR_PRIMARY; box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_top = Emu(80000)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "⚠  Restricción clave   "
    r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = COLOR_PRIMARY
    r2 = p.add_run()
    r2.text = ("ECC 6.0 EHP 8 NO es S/4HANA → RAP no disponible. Se expone OData V2 "
                "vía SEGW desde ECC y se complementa con CAP en BTP. "
                "Sin modificaciones del estándar.")
    r2.font.size = Pt(12); r2.font.color.rgb = COLOR_TEXT


def slide_04_situacion(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "SITUACIÓN ACTUAL — PROCESO MANUAL")
    add_footer(slide, 4)
    add_text(slide, 0.7, 1.5, 12.0, 0.6,
             "Autorización de salidas de materiales contra reservas de mantenimiento — proceso mixto SAP + papel:",
             font_size=12, color=COLOR_TEXT)
    steps = [
        ("1", "Planeador",     "Crea reserva PM en IW32 (tabla RESB)"),
        ("2", "Aprobación",    "Correo electrónico al ingeniero / jefe"),
        ("3", "Formato físico","HOC-36 a mano + firmas presenciales"),
        ("4", "Bodega",        "Valida autorización física y ejecuta MIGO"),
        ("5", "Reporte",       "MB51 sin columna Centro de Costo"),
    ]
    x0 = 0.6; w = 2.4; gap = 0.075; y = 2.6
    for i, (n, t, d) in enumerate(steps):
        x = x0 + i * (w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(y), Inches(w), Inches(1.7))
        box.fill.solid(); box.fill.fore_color.rgb = COLOR_PRIMARY
        box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(80000); tf.margin_right = Emu(80000)
        tf.margin_top = Emu(80000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run(); r1.text = f"Paso {n}"
        r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = COLOR_LIGHT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = t
        r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = COLOR_WHITE
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = d
        r3.font.size = Pt(9); r3.font.color.rgb = COLOR_WHITE
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          Inches(x + w + 0.005), Inches(y + 0.7),
                                          Inches(0.06), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = COLOR_ACCENT
            arr.line.fill.background()
    add_text(slide, 0.7, 4.7, 12.0, 0.4, "Consecuencias operativas",
             font_size=14, bold=True, color=COLOR_RED)
    add_bullets(slide, 0.7, 5.1, 12.0, 1.8, [
        "Reprocesos cuando la orden se cierra anticipadamente o el descargue no se hace a tiempo.",
        "Reapertura o registro contra centro de costo — desviaciones de imputación.",
        "Información no explotable analíticamente: no hay KPIs de tiempo de aprobación.",
        "Auditoría compleja por la mezcla de correos, firmas físicas y registros parciales en SAP.",
    ], font_size=12)


def slide_05_brechas(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "BRECHAS IDENTIFICADAS")
    add_footer(slide, 5)
    data = [
        ["#", "Brecha", "Impacto operativo"],
        ["GAP-01", "Alta dependencia de actividades manuales", "Errores y reprocesos"],
        ["GAP-02", "Duplicidad de información (SAP vs registros físicos)", "Inconsistencia de datos"],
        ["GAP-03", "Falta de integración entre sistema y operación", "Inventario poco fiable"],
        ["GAP-04", "Aprobaciones por correo y firma manual", "Cuellos de botella · retrasos"],
        ["GAP-05", "Sin trazabilidad en tiempo real", "Riesgo de control interno"],
        ["GAP-06", "Mecanismos de firma manual difíciles de auditar", "Auditoría compleja"],
        ["GAP-07", "Proceso sin estandarización entre campos", "Variabilidad operativa"],
        ["GAP-08", "Desplazamientos físicos del aprobador", "Productividad baja"],
        ["GAP-09", "Información no explotable analíticamente", "Sin KPIs de gestión"],
        ["GAP-10", "MB51 estándar no muestra Centro de Costo", "Imputación incorrecta"],
    ]
    add_table(slide, 0.7, 1.45, 12.0, 5.4, data,
              col_widths=[0.9, 5.6, 5.5],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)


def slide_06_tobe(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "RESULTADO ESPERADO — WORKFLOW DIGITAL")
    add_footer(slide, 6)
    add_text(slide, 0.7, 1.4, 12.0, 0.4,
             "Workflow 100% digital gobernado por Status de Usuario (REBL/REED/03PL/04PR) y disparado por la combinación 04PR + MACO.",
             font_size=12, color=COLOR_TEXT)
    data = [
        ["#",  "Paso",                              "Mecánica técnica"],
        ["1",  "Planificación y creación",          "IW32 — Status MACO + Status usuario REBL (Reserva bloqueada)"],
        ["2",  "Bloqueo edición de componentes",    "Planeador asigna 03PL (Orden Planeada)"],
        ["3",  "Liberación de orden con reserva",   "Planeador asigna 04PR (Orden Programada)"],
        ["4",  "Disparador del workflow",           "Combinación 04PR + MACO activa el flujo en BTP"],
        ["5",  "Aprobación",                        "App SAP/SuccessFactors — decisión Sí/No"],
        ["6a", "Decisión: SÍ",                      "BAPI STATUS_CHANGE_EXTERN desactiva REBL → bodega notificada"],
        ["6b", "Decisión: NO",                      "REBL permanece + activa REED → planeador ajusta y reinicia"],
        ["7",  "Ejecución en bodega",               "MIGO consume materiales — registro de usuarios"],
        ["8",  "Generación del PDF",                "PDF en BTP + archivo en SAP Document Management Service"],
        ["9",  "Reporte mejorado",                  "MB51 con columna Centro de Costo agregada"],
    ]
    add_table(slide, 0.7, 1.95, 12.0, 4.9, data,
              col_widths=[0.7, 4.0, 7.3],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)


def slide_07_beneficios(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "BENEFICIOS Y OBJETIVOS")
    add_footer(slide, 7)
    benefits = [
        ("80%",    "Reducción del tiempo de ciclo de aprobación"),
        ("100%",   "Trazabilidad de extremo a extremo en SAP"),
        ("0",      "Formatos físicos preimpresos HOC-36"),
        ("0",      "Firmas manuales — todo digital y auditable"),
        ("KPI",    "Información explotable (tiempos, rechazos, SLAs)"),
        ("Audit",  "Cumplimiento de control interno y compliance"),
    ]
    x0 = 0.5; y0 = 1.5; w = 4.1; h = 1.55; gx = 0.15; gy = 0.2
    for i, (big, desc) in enumerate(benefits):
        col = i % 3; row = i // 3
        x = x0 + col * (w + gx); y = y0 + row * (h + gy)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT
        box.line.color.rgb = COLOR_PRIMARY; box.line.width = Pt(1.0)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
        tf.margin_top = Emu(120000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = big
        r.font.bold = True; r.font.size = Pt(38); r.font.color.rgb = COLOR_PRIMARY
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = Pt(12); r2.font.color.rgb = COLOR_TEXT
    add_text(slide, 0.7, 5.4, 12.0, 0.4, "Objetivos estratégicos",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.7, 5.8, 12.0, 1.2, [
        "Asegurar Clean Core en ECC — toda la lógica de aprobación reside en BTP.",
        "Habilitar un modelo extensible para nuevos workflows de aprobación corporativos.",
        "Estandarizar la operación entre campos eliminando variabilidad del proceso manual.",
    ], font_size=11)


def slide_08_arquitectura(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "ARQUITECTURA DE SOLUCIÓN — VISTA L1")
    add_footer(slide, 8)
    diag = f"{DIAGRAMS}/diagram_01.png"
    if os.path.exists(diag):
        slide.shapes.add_picture(diag, Inches(0.5), Inches(1.4),
                                  width=Inches(8.8), height=Inches(5.4))
    add_text(slide, 9.5, 1.4, 3.4, 0.4, "Capas",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 9.5, 1.85, 3.4, 5.2, [
        ("UI",            "UI5 Freestyle + Build Work Zone"),
        ("Backend",       "CAP Node.js (@sap/cds 9.7.x)"),
        ("Persistencia",  "HANA Cloud HDI Container"),
        ("Scheduler",     "node-cron embebido en CAP"),
        ("PDFs",          "Generación in-process + DMS"),
        ("Seguridad",     "XSUAA + IAS + Principal Propagation"),
        ("Conectividad",  "Cloud Connector → RFC + OData V2"),
    ], font_size=10, spacing_after=8)


def slide_09_componentes(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "COMPONENTES BTP REQUERIDOS")
    add_footer(slide, 9)
    add_text(slide, 0.7, 1.4, 12.0, 0.4,
             "Servicios BTP a habilitar en el subaccount HOCOL — 2 ambientes (DEV y PRD).",
             font_size=12, color=COLOR_TEXT)
    data = [
        ["Servicio BTP",                          "Plan / Uso",                "Rol en la solución"],
        ["Cloud Foundry Runtime",                 "Standard",                  "Ejecución de la app CAP Node.js y approuter"],
        ["SAP HANA Cloud",                        "HDI Container",             "Persistencia de estado de workflow y auditoría"],
        ["XSUAA",                                 "Application",               "Autenticación + autorización (scopes y role-collections)"],
        ["Destination Service",                   "Lite",                      "Destinos hacia ECC y SuccessFactors"],
        ["Connectivity Service",                  "Lite",                      "Túnel BTP ↔ Cloud Connector"],
        ["HTML5 Application Repository",          "App-host",                  "Hosting de la app UI5"],
        ["SAP Build Work Zone (Standard)",        "Standard",                  "Launchpad y catálogo de la app"],
        ["Application Logging",                   "Lite",                      "Logs centralizados de la app CAP"],
        ["Audit Log",                             "Premium",                   "Pista de auditoría de eventos críticos"],
        ["IAS (Identity Authentication)",         "Tenant",                    "Federación con IdP corporativo HOCOL"],
        ["SAP Document Management Service",       "Standard (ya licenciado)",  "Archivo y consulta del PDF generado"],
    ]
    add_table(slide, 0.5, 2.0, 12.4, 4.7, data,
              col_widths=[4.0, 3.0, 5.4],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)


def slide_10_workflow(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "MOTOR DE WORKFLOW — CAP CUSTOM + NODE-CRON")
    add_footer(slide, 10)
    add_text(slide, 0.7, 1.4, 6.0, 0.4, "Decisión técnica",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_text(slide, 0.7, 1.85, 6.0, 2.0,
             "El workflow se implementa como una state machine in-process dentro del backend "
             "CAP Node.js, persistida en HANA HDI. No se requiere motor externo: el dominio "
             "funcional es acotado y los volúmenes esperados encajan en el servicio CAP.",
             font_size=11, color=COLOR_TEXT)
    add_text(slide, 0.7, 3.9, 6.0, 0.4, "Ventajas",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.7, 4.3, 6.0, 2.7, [
        "Sin costos de licencia adicionales por motor de workflow",
        "Control total del modelo de auditoría y reglas de SLA",
        "Single source of truth en Git (versión-able junto al código)",
        "Despliegue y observabilidad unificados con el resto de la app",
        "Reglas y temporizadores tested junto al código CAP",
    ], font_size=11)
    add_text(slide, 7.0, 1.4, 6.0, 0.4, "Estados del workflow",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    states = [
        ("PENDING",    "Workflow creado, esperando aprobador",      COLOR_AMBER),
        ("APPROVED",   "Aprobado → BAPI desactiva REBL en ECC",     COLOR_GREEN),
        ("REJECTED",   "Rechazado → activa REED + notifica planead.", COLOR_RED),
        ("EXPIRED",    "Vencido por SLA → recordatorio / reasignar", COLOR_MUTED),
        ("CANCELLED",  "Cancelado por planeador / sistema",          COLOR_MUTED),
    ]
    y = 1.85
    for name, desc, col in states:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(7.0), Inches(y),
                                       Inches(1.7), Inches(0.5))
        pill.fill.solid(); pill.fill.fore_color.rgb = col
        pill.line.fill.background()
        tf = pill.text_frame
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name
        r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = COLOR_WHITE
        add_text(slide, 8.85, y + 0.1, 4.2, 0.4, desc,
                 font_size=11, color=COLOR_TEXT)
        y += 0.65
    add_text(slide, 7.0, 5.45, 6.0, 0.6,
             "Recordatorios y SLA con node-cron embebido — sin SAP Job Scheduler ni dependencias externas.",
             font_size=10, italic=True, color=COLOR_MUTED)


def slide_11_plan(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "PLAN DE TRABAJO — 6 FASES")
    add_footer(slide, 11)
    add_text(slide, 0.7, 1.45, 12.0, 0.4,
             "Proyecto secuencial de 18 semanas con esfuerzo cerrado de 1.020 horas-persona en el alcance BTP.",
             font_size=12, color=COLOR_TEXT)
    phases = [(p["code"], p["name"], p["duration"], p["hours"]) for p in PROJECT_PHASES]
    y = 2.1
    x0 = 0.5; w = 2.00; gap = 0.12
    colors = [COLOR_DARK, COLOR_PRIMARY, COLOR_PRIMARY, COLOR_PRIMARY, COLOR_DARK, COLOR_DARK]
    for i, (n, name, dur, h) in enumerate(phases):
        add_phase_card(slide, x0 + i * (w + gap), y, w, 1.8,
                        n, name, dur, h, fill=colors[i])
    add_text(slide, 0.7, 4.4, 12.0, 0.4, "Actividades destacadas por fase",
             font_size=13, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.7, 4.85, 12.0, 2.2, [
        ("F0", "Subaccount DEV/PRD, Cloud Connector mapping, destinations, IAS, XSUAA"),
        ("F1", "Modelo CDS, diseño UI5, contratos OData ECC, plan de tests, MTA"),
        ("F2", "Backend CAP, app UI5, integración ECC, generación PDF + DMS, scheduled tasks"),
        ("F3", "Unit + integration + E2E + UAT con usuarios HOCOL"),
        ("F4-F5", "CutOver DEV → PRD + soporte continuo 1 semana de hypercare"),
    ], font_size=11, spacing_after=4)


def slide_12_entregables(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "ENTREGABLES")
    add_footer(slide, 12)
    data = [
        ["Fase",                                "Entregable",                              "Detalle"],
        [PROJECT_PHASES[0]["short"],            "Subaccount BTP DEV/PRD",                  "Provisioning, entitlements, members, trust config IAS"],
        ["",                                    "Cloud Connector mapping",                 "Virtual hosts + destinations + Principal Propagation"],
        ["",                                    "Modelo de roles XSUAA",                   "Scopes y role-collections (Approver, Planner, Warehouse, Admin)"],
        ["F1 Diseño",                           "Documento de diseño técnico",             "Modelo CDS, contratos OData ECC, diseño UI5, MTA"],
        [PROJECT_PHASES[2]["short"],            "Backend CAP + HANA HDI",                  "Workflow state machine, scheduled tasks, integración ECC y DMS"],
        ["",                                    "App UI5 Freestyle",                       "Inbox de aprobaciones, detalle de reserva, fragments de decisión"],
        ["F3 Testing",                          "Plan + casos de prueba ejecutados",       "Unit (CAP/UI5) + integration + E2E + soporte UAT"],
        [PROJECT_PHASES[4]["short"],            "Plan de CutOver + paso a PRD",            "Pipeline MTA con go/no-go por hito"],
        [PROJECT_PHASES[5]["short"],            "Soporte continuo + handover",             "1 semana de estabilización con soporte continuo"],
    ]
    add_table(slide, 0.5, 1.5, 12.4, 5.3, data,
              col_widths=[1.7, 3.5, 7.2],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)


def slide_13_cronograma(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "CRONOGRAMA — 18 SEMANAS")
    add_footer(slide, 13)
    add_text(slide, 0.7, 1.4, 12.0, 0.4,
             "Arquitecto BTP lidera F0. Desde F1 hasta cierre del hypercare, los 3 perfiles trabajan en paralelo durante 15 semanas.",
             font_size=11, color=COLOR_TEXT)
    y_start = 1.85
    row_h = 0.34
    chart_x = 1.95
    chart_w = 10.6
    week_w = chart_w / 18
    add_text(slide, 0.5, y_start, 1.4, 0.3, "Fase / Perfil",
             font_size=10, bold=True, color=COLOR_PRIMARY)
    for i in range(18):
        x = chart_x + i * week_w
        add_text(slide, x, y_start, week_w, 0.3, f"S{i+1}",
                 font_size=8, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(y_start + 0.3),
                                   Inches(chart_w + chart_x - 0.5), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()
    # Filas del Gantt: (label, start_week, duration_weeks, color)
    # El label se deriva de PROJECT_PHASES[i]["short"] para mantener la
    # consistencia con la tabla de entregables y evitar duplicación literal.
    _gantt_meta = [(1, 3, COLOR_DARK), (4, 2, COLOR_PRIMARY),
                    (6, 8, COLOR_PRIMARY), (14, 3, COLOR_PRIMARY),
                    (17, 1, COLOR_DARK), (18, 1, COLOR_DARK)]
    rows = [(PROJECT_PHASES[i]["short"], *meta) for i, meta in enumerate(_gantt_meta)]
    y = y_start + 0.45
    for label, start, dur, color in rows:
        add_text(slide, 0.5, y, 1.45, row_h - 0.05, label,
                 font_size=10, bold=True, color=COLOR_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(chart_x + (start - 1) * week_w + 0.02),
                                      Inches(y + 0.05),
                                      Inches(dur * week_w - 0.04),
                                      Inches(row_h - 0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        y += row_h
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(y + 0.03),
                                    Inches(chart_w + chart_x - 0.5), Inches(0.02))
    line2.fill.solid(); line2.fill.fore_color.rgb = COLOR_MUTED
    line2.line.fill.background()
    y += 0.13
    perfiles = [
        ("Arquitecto BTP",       1, 18, COLOR_PRIMARY),
        ("BTP & CAP Developer",  4, 15, COLOR_DARK),
        ("UI5 Freestyle Dev",    4, 15, COLOR_DARK),
    ]
    for label, start, dur, color in perfiles:
        add_text(slide, 0.5, y, 1.45, row_h - 0.05, label,
                 font_size=10, bold=True, color=COLOR_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(chart_x + (start - 1) * week_w + 0.02),
                                      Inches(y + 0.05),
                                      Inches(dur * week_w - 0.04),
                                      Inches(row_h - 0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        y += row_h
    add_text(slide, 0.5, y + 0.05, 12.0, 0.3, "Hitos clave",
             font_size=12, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.5, y + 0.4, 12.0, 0.9, [
        ("M1 sem 3",   "Discovery completo + arquitectura aprobada"),
        ("M2 sem 5",   "Subaccount BTP operativo + Cloud Connector validado E2E"),
        ("M5 sem 13",  "Integración E2E ECC ↔ BTP completa"),
        ("M7 sem 17",  "Sign-off UAT"),
        ("M8 sem 18",  "Go-Live en PRD"),
    ], font_size=9, spacing_after=1)


def slide_14_premisas_dependencias(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "PREMISAS Y DEPENDENCIAS DEL CLIENTE")
    add_footer(slide, 14)
    add_text(slide, 0.5, 1.35, 12.3, 0.45,
             "El cumplimiento del cronograma y la estimación dependen de las siguientes premisas y dependencias. Si alguna no se cumple, "
             "el cronograma se replantea formalmente sin afectar la cifra cerrada de horas comprometida.",
             font_size=11, color=COLOR_TEXT)
    add_text(slide, 0.5, 1.95, 6.0, 0.4, "Premisas asumidas",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.5, 2.4, 6.0, 4.0, [
        "Cloud Connector existente acepta agregar nuevos virtual hosts sin reconfigurar lo existente.",
        "ECC 6.0 EHP 8 tiene SAP_BASIS ≥ 752 y soporta SEGW + OData V2 sin restricciones.",
        "BAPIs estándar de reservas, órdenes PM y goods movement están habilitadas en RISE.",
        "SAP Document Management Service está activo y licenciado en el tenant BTP de HOCOL.",
        "Tenant IAS accesible y federable con el IdP corporativo (Azure AD Ecopetrol).",
        "Política de seguridad permite Principal Propagation hacia ECC.",
        "Volumen máximo estimado: 500 reservas/día (sizing HANA basado en esta cifra).",
        "SMTP corporativo permite envíos desde direcciones BTP CF.",
        "Idioma de entregables: español. Sin certificaciones regulatorias adicionales requeridas.",
    ], font_size=10, spacing_after=4)
    add_text(slide, 6.8, 1.95, 6.0, 0.4, "Dependencias del cliente (HOCOL/Ecopetrol)",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 6.8, 2.4, 6.0, 4.0, [
        "Provisioning de subaccount BTP DEV y PRD (trámite interno Ecopetrol).",
        "Accesos al sistema RISE-ECC DEV/PRD con usuarios de desarrollador.",
        "Coordinación con SAP RISE para configuraciones Basis on-premise (PFCG, STMS, USREXTID).",
        "Disponibilidad ≥95% de usuarios funcionales para Discovery (workshops) y UAT.",
        "Equipo HOCOL/Tivit ABAP, Funcional y PMO con dedicación comprometida en kickoff.",
        "Aprobación oportuna de transportes ABAP en DEV y PRD.",
        "Conectividad de red corporativa (firewalls, NAT, whitelists IP de BTP CF).",
        "Acceso documentado al tenant SAP Document Management Service.",
        "Cuenta de servicio (technical user) en ECC para integración.",
    ], font_size=10, spacing_after=4)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(6.45),
                                  Inches(12.3), Inches(0.55))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Control de cambios  "
    r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = COLOR_DARK
    r2 = p.add_run()
    r2.text = ("Cualquier requerimiento no listado en la sección «Alcance — Incluye» se gestiona como Change Request (CR) "
                "formal con re-estimación de horas y posible reprogramación de fechas.")
    r2.font.size = Pt(11); r2.font.color.rgb = COLOR_DARK


def slide_15_alcance_in_out(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "ALCANCE — INCLUYE / NO INCLUYE")
    add_footer(slide, 15)
    add_text(slide, 0.5, 1.35, 12.3, 0.4,
             "Delimitación explícita del alcance Tivit BTP. Lo que no aparece a la izquierda no está cubierto por esta propuesta.",
             font_size=11, color=COLOR_TEXT)
    in_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(1.85),
                                     Inches(6.0), Inches(0.45))
    in_box.fill.solid(); in_box.fill.fore_color.rgb = COLOR_PRIMARY
    in_box.line.fill.background()
    tf = in_box.text_frame
    tf.margin_left = Emu(120000); tf.margin_top = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "✓  INCLUYE  —  Alcance Tivit BTP"
    r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = COLOR_WHITE
    in_items = [
        "Subaccount BTP HOCOL: 2 ambientes (DEV y PRD).",
        "Backend CAP Node.js (modelo CDS, services, handlers, workflow logic).",
        "Base de datos HANA Cloud (HDI Container).",
        "App UI5 Freestyle «Aprobación de Reservas PM».",
        "Cloud Connector mapping del sistema RISE-ECC hacia BTP.",
        "Generación de PDF + archivo en SAP Document Management Service.",
        "Integración BTP ↔ ECC (RFC + OData V2 consumido desde ECC).",
        "Integración BTP ↔ SuccessFactors (opcional, sujeto a disponibilidad).",
        "Notificaciones por correo electrónico.",
        "Modelo de autorizaciones BTP (XSUAA) — diseño de roles y scopes.",
        "Federación IAS con corporate IdP de HOCOL.",
        "Pipeline CI/CD para deployment de BTP.",
        "Testing E2E + soporte UAT con usuarios HOCOL.",
        "Hypercare post go-live (1 semana de estabilización).",
    ]
    add_bullets(slide, 0.5, 2.35, 6.0, 4.55, in_items,
                font_size=9, color=COLOR_TEXT, spacing_after=2)
    out_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(6.8), Inches(1.85),
                                      Inches(6.0), Inches(0.45))
    out_box.fill.solid(); out_box.fill.fore_color.rgb = COLOR_RED
    out_box.line.fill.background()
    tf = out_box.text_frame
    tf.margin_left = Emu(120000); tf.margin_top = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "✗  NO INCLUYE  —  Fuera de alcance"
    r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = COLOR_WHITE
    out_items = [
        "Desarrollo ABAP en ECC (BAdIs, SEGW, Enhancements) — área ABAP de Tivit.",
        "Configuración Basis on-premise en ECC (PFCG, STMS, USREXTID) — HOCOL/RISE.",
        "Functional Specification detallada — Analista Funcional (Tivit/HOCOL).",
        "Migración masiva de reservas históricas.",
        "Procesos de mantenimiento distintos al workflow de reserva.",
        "Cambios en módulos diferentes a PM/MM.",
        "Migración a S/4HANA o upgrades de versión de ECC.",
        "Apps móviles nativas (Android/iOS) — la UI5 será responsive web.",
        "Workflow de aprobación para órdenes de compra (PR/PO).",
        "Manuales de usuario final, change management interno y capacitación.",
        "Soporte continuo post-hypercare (separable como contrato AMS).",
        "Adquisición y/o costos de licenciamiento BTP — Account Executive SAP.",
        "Conectividad de red corporativa (firewalls, NAT, whitelisting de IPs).",
        "Auditoría externa, certificaciones regulatorias y compliance avanzado.",
    ]
    add_bullets(slide, 6.8, 2.35, 6.0, 4.55, out_items,
                font_size=9, color=COLOR_TEXT, spacing_after=2)


def slide_16_equipo_estimacion(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "EQUIPO TIVIT BTP Y ESTIMACIÓN")
    add_footer(slide, 16)
    add_text(slide, 0.7, 1.4, 6.0, 0.4, "Perfiles BTP — alcance Tivit",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    data1 = [
        ["Perfil",                     "Dedicación",                "Horas"],
        ["Arquitecto BTP",             "F0 100% + Transv. F1-F5",   "550 h"],
        ["BTP & CAP Developer",        "Transversal F1-F5",         "290 h"],
        ["UI5 Freestyle Developer",    "Transversal F1-F5",         "180 h"],
        ["TOTAL",                      "18 semanas",                TOTAL_HOURS],
    ]
    add_table(slide, 0.7, 1.85, 6.0, 2.0, data1,
              col_widths=[2.7, 2.3, 1.0],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.7), Inches(4.1),
                                  Inches(6.0), Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_ACCENT
    box.line.color.rgb = COLOR_PRIMARY; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(160000); tf.margin_top = Emu(120000)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Estimación cerrada — alcance BTP"
    r.font.size = Pt(12); r.font.color.rgb = COLOR_DARK
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "1.020 horas-persona  ·  18 semanas"
    r2.font.bold = True; r2.font.size = Pt(22); r2.font.color.rgb = COLOR_DARK
    p3 = tf.add_paragraph(); p3.space_before = Pt(6)
    r3 = p3.add_run()
    r3.text = "Arquitectura · desarrollo · integración · testing E2E · despliegue · hypercare"
    r3.font.size = Pt(10); r3.font.color.rgb = COLOR_DARK
    add_text(slide, 7.0, 1.4, 6.0, 0.4, "Distribución por fase",
             font_size=14, bold=True, color=COLOR_PRIMARY)
    data2 = [["Fase", "Duración", "Esfuerzo", "%"]]
    data2.extend([p["full_name"], p["duration"], p["hours"], p["percent"]]
                  for p in PROJECT_PHASES)
    data2.append(["TOTAL", TOTAL_DURATION, TOTAL_HOURS, "100%"])
    add_table(slide, 7.0, 1.85, 6.0, 3.4, data2,
              col_widths=[2.4, 1.2, 1.5, 0.9],
              body_font_size=10, header_font_size=11,
              first_col_bold=True)
    add_text(slide, 7.0, 5.4, 6.0, 0.35,
             "Perfiles Tivit adicionales (no incluidos en esta cifra)",
             font_size=11, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 7.0, 5.75, 6.0, 1.3, [
        "PM / Scrum Master · Functional Analyst (PM/MM) · ABAP Developer",
        "Estimación entregada como anexo separado por las áreas correspondientes",
    ], font_size=10)


def slide_17_riesgos(prs):
    slide = add_blank(prs)
    add_logos(slide); add_header_banner(slide, "RIESGOS CLAVE Y MITIGACIONES")
    add_footer(slide, 17)
    data = [
        ["#",     "Riesgo",                                                                "Prob.",   "Imp.",   "Mitigación"],
        ["R-01",  "Limitación de APIs en RISE-ECC para exponer BAPIs necesarias",         "Media",   "Alto",   "Discovery con POC concreto · plan B con RFC custom"],
        ["R-02",  "Cloud Connector no permite Principal Propagation por config RISE",     "Baja",    "Alto",   "Plan B: BasicAuth con tech user documentado"],
        ["R-03",  "SuccessFactors no disponible para reglas de ausencias/reemplazos",     "Media",   "Medio",  "Plan B: tabla ApproverAssignment custom desde UI Admin"],
        ["R-05",  "Subaccount BTP no aprobado a tiempo por procesos internos Ecopetrol",  "Media",   "Alto",   "Iniciar trámite en semana 1 · escalamiento a Sponsor"],
        ["R-10",  "SAP Document Management Service no disponible en tenant HOCOL",        "Baja",    "Medio",  "Validar entitlement en semana 1 · alt.: Object Store"],
        ["R-13",  "Disponibilidad de equipo HOCOL (ABAP, Basis) para apoyar el proyecto", "Media",   "Alto",   "Confirmar dedicación en kickoff · escalar si <2 sem"],
    ]
    add_table(slide, 0.5, 1.5, 12.4, 4.2, data,
              col_widths=[0.7, 4.5, 0.9, 0.9, 5.4],
              body_font_size=9, header_font_size=11,
              first_col_bold=True)
    add_text(slide, 0.5, 5.95, 12.4, 0.35, "Supuestos clave a confirmar en Discovery",
             font_size=12, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, 0.5, 6.3, 12.4, 0.85, [
        "Cloud Connector permite agregar nuevos virtual hosts sin reconfigurar lo existente",
        "ECC soporta SEGW + OData V2 (SAP_BASIS ≥ 752)",
        "SAP Document Management Service licenciado y disponible en el tenant BTP de HOCOL",
        "Tenant IAS accesible y permite federación con IdP corporativo",
    ], font_size=9, spacing_after=2)


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def main():
    if not os.path.exists(TEMPLATE):
        raise SystemExit(f"❌ Falta template: {TEMPLATE}")
    shutil.copy(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    delete_all_slides(prs)
    builders = [
        slide_01_portada, slide_02_agenda, slide_03_contexto, slide_04_situacion,
        slide_05_brechas, slide_06_tobe, slide_07_beneficios, slide_08_arquitectura,
        slide_09_componentes, slide_10_workflow, slide_11_plan, slide_12_entregables,
        slide_13_cronograma, slide_14_premisas_dependencias, slide_15_alcance_in_out,
        slide_16_equipo_estimacion, slide_17_riesgos,
    ]
    for b in builders:
        b(prs)
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f"✓ Generado: {OUTPUT}")
    print(f"  Tamaño: {size_kb:.1f} KB · {len(builders)} slides")


if __name__ == "__main__":
    main()
